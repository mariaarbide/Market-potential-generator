# ==============================================================================
# SCRIPT COMPLETO Y FINAL DE LA APLICACIÓN DE MARKET INTELLIGENCE
# Restaura la lógica de análisis detallada del usuario y la une con la
# funcionalidad de exportación a PDF y PPTX.
#
# REQUISITOS:
# 1. Instalar librerías:
#    pip install streamlit pandas requests fpdf2 python-pptx
#
# 2. Archivos en la misma carpeta que este script:
#    - vicomtech_logo.png
#    - vicomtech_template.pptx
#
# 3. Configurar Claves de API:
#    - GOOGLE_API_KEY, GOOGLE_CSE_ID, GEMINI_API_KEY
# ==============================================================================

import streamlit as st
import requests
import re
import pandas as pd
import urllib3
import json
import time
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import io

# --- 1. CONFIGURACIÓN INICIAL ---

# Desactivar warnings de SSL
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# --- Claves de API (reemplaza con las tuyas) ---
GOOGLE_API_KEY = "AIzaSyBLliyWr_A_F4_-E7tr3G5pWurlZejcCfA"
GOOGLE_CSE_ID = "244bd67fb4e7a4b81"
GEMINI_API_KEY = "AIzaSyC5xpVqN_IuCiLLgb737Ca2Ep6cINmDvM8"

# --- Modo de Simulación y Datos Mock ---
MOCK_GOOGLE_SEARCH_RESPONSES = False
# Datos de simulación completos
MOCK_DATA = {
    # Mock data para "Healthcare Chatbots"
    "Healthcare Chatbots site:statista.com": { "items": [ {"title": "Healthcare chatbots market size worldwide 2023-2030 | Statista", "link": "https://www.statista.com/statistics/1083438/healthcare-chatbots-market-size-worldwide/", "snippet": "Global healthcare chatbot market analysis. Data from Statista on market size trends. Market size of healthcare chatbots was valued at USD 1.20 billion in 2024 and is projected to reach USD 10.26 billion by 2034, with a CAGR of 23.92%."} ] },
    "Healthcare Chatbots site:marketdataforecast.com": { "items": [ {"title": "Healthcare Chatbot Market Size, Share & Growth Report (2024-2029)", "link": "https://www.marketdataforecast.com/market-reports/healthcare-chatbots-market", "snippet": "Comprehensive report on healthcare chatbot market. Forecasts and growth. The market is expected to grow at a CAGR of 21.5% from 2024 to 2029."} ] },
    "Healthcare Chatbots site:precedenceresearch.com": { "items": [ {"title": "Healthcare Chatbots Market Size, Share, Trends, & Forecast (2024-2034)", "link": "https://www.precedenceresearch.com/healthcare-chatbots-market", "snippet": "Precedence Research insights into the global healthcare chatbots market, forecasts up to 2034. The global healthcare chatbots market size was valued at USD 1.20 billion in 2024 and is projected to hit around USD 10.26 billion by 2034 with a CAGR of 23.92%."} ] },
    "Healthcare Chatbots market size in USD total addressable serviceable obtainable": { "items": [ {"snippet": "The global healthcare chatbots market size was valued at USD 1.20 billion in 2024 and is projected to hit around USD 10.26 billion by 2034 with a CAGR of 23.92%. This represents the Total Addressable Market (TAM). The Serviceable Available Market (SAM) for digital pediatric solutions is estimated at USD 370 million, with a Serviceable Obtainable Market (SOM) in Spain of USD 0.5-1.2 million.", "title": "Healthcare Chatbots Market Overview", "link": "https://example.com/market-overview"} ] },
    "Healthcare Chatbots competitors OR similar companies crunchbase.com": { "items": [ {"title": "Ada Health - Crunchbase Company Profile & Funding", "link": "https://www.crunchbase.com/organization/ada-health", "snippet": "Ada Health is a global health company that provides AI-powered health assessment and care navigation solutions. The company has raised over $300M in funding. Headquartered in Germany. Offers AI-powered symptom assessment for general users."} ] },
    "Healthcare Chatbots market drivers and barriers": { "items": [ {"snippet": "Drivers for healthcare chatbots include increasing demand for remote care, advancements in AI, and cost reduction. Barriers include regulatory hurdles, data privacy concerns, and user trust issues.", "title": "Trends in Healthcare Chatbots", "link": "https://example.com/drivers"} ] },
    "Healthcare Chatbots growth factors challenges": { "items": [ {"snippet": "Growth factors for healthcare chatbots are improved accessibility and efficiency. Challenges involve integration with existing systems and maintaining accuracy for complex medical conditions.", "title": "Key Challenges and Growth Drivers", "link": "https://example.com/challenges"} ] },
    "Healthcare Chatbots industry trends opportunities threats": { "items": [ {"snippet": "Major trends include personalized AI and predictive analytics. Opportunities lie in chronic disease management and mental health support. Threats involve cybersecurity risks and competition from traditional healthcare providers.", "title": "Future of Healthcare AI", "link": "https://example.com/trends"} ] },
    "Healthcare Chatbots M&A collaborations OR acquisitions OR partnerships site:crunchbase.com": { "items": [ {"title": "Amwell Acquires Conversa Health - Crunchbase News", "link": "https://techcrunch.com/amwell-conversa-acquisition", "snippet": "2020: Amwell acquires Conversa Health (health chatbot) for USD 320M, reinforcing its suite with conversational AI and SilverCloud. Type: Strategic acquisition."} ] }
}


# --- 2. CONFIGURACIÓN DE LA PÁGINA Y ESTILOS ---
st.set_page_config(page_title="Market Intelligence Vicomtech", page_icon="vicomtech_logo.png", layout="centered")

st.markdown("""
    <style>
    .main { background-color: #f0f4f8; }
    .stButton>button { background-color: #0077c2; color: white; border-radius: 5px; padding: 8px 18px; font-weight: 600; }
    .stButton>button:hover { background-color: #005a99; }
    .stTextInput>div>div>input { border: 2px solid #0077c2; border-radius: 5px; }
    header, footer { display: none; }
    h1, h2, h3 { color: #0077c2; }
    </style>
""", unsafe_allow_html=True)

# --- 3. FUNCIONES DETALLADAS DE ANÁLISIS Y EXTRACCIÓN ---

@st.cache_data(show_spinner=False)
def run_full_analysis(query):
    """
    Función que ejecuta todo el análisis detallado y devuelve los resultados.
    """
    results = {}
    # Esta función ahora agrupa todas las llamadas de análisis
    results['statista_reports'] = buscar_en_fuente(query, "statista.com")
    results['market_data_reports'] = buscar_en_fuente(query, "marketdataforecast.com")
    results['precedence_reports'] = buscar_en_fuente(query, "precedenceresearch.com")
    
    cagr, market_values, tam, sam, som, snippet = obtener_numeros_cruciales(query)
    results['cagr'] = cagr
    results['market_values'] = market_values
    results['tam'] = tam
    results['sam'] = sam
    results['som'] = som
    results['tam_sam_som_snippet'] = snippet
    
    results['empresas_competidoras'] = extraer_empresas_relacionadas(query)
    results['ma_collaborations'] = extraer_colaboraciones_ma(query)
    results['drivers'], results['barriers'] = generar_drivers_barreras(query)

    return results

def safe_request(url, params=None):
    try:
        resp = requests.get(url, params=params, verify=False, timeout=15)
        resp.raise_for_status()
        return resp.json()
    except requests.exceptions.RequestException as e:
        st.error(f"Error en la petición HTTP: {e}.")
        return None

def search_google_custom_search(query_text):
    if MOCK_GOOGLE_SEARCH_RESPONSES:
        return MOCK_DATA.get(query_text, {"items": []})
    if not GOOGLE_API_KEY or not GOOGLE_CSE_ID:
        st.error("Por favor, configura GOOGLE_API_KEY y GOOGLE_CSE_ID.")
        return {"items": []}
    params = {"key": GOOGLE_API_KEY, "cx": GOOGLE_CSE_ID, "q": query_text, "num": 10}
    url = "https://www.googleapis.com/customsearch/v1"
    data = safe_request(url, params)
    time.sleep(1)
    return data if data else {"items": []}

def buscar_en_fuente(query_text, site):
    resultados_raw = search_google_custom_search(f"{query_text} site:{site}")
    found_results = []
    for r in resultados_raw.get("items", [])[:3]:
        if r.get("title") and r.get("link"):
            found_results.append((r["title"], r["link"]))
    return found_results

def convertir_valor(num_str, unidad):
    if num_str is None: return None
    num_str = num_str.replace("$", "").replace(",", "").strip()
    try: num = float(num_str)
    except ValueError: return None
    unidad = unidad.lower() if unidad else ""
    if unidad in ["billion", "bn", "billones"]: return round(num * 1000, 2)
    if unidad in ["million", "m", "millones"]: return round(num, 2)
    return round(num, 2)

def obtener_numeros_cruciales(product_query):
    busqueda = f"{product_query} market size growth CAGR annual value"
    results_raw = search_google_custom_search(busqueda)
    results = results_raw.get("items", [])
    cagr_found, market_values = None, []
    for r in results:
        snippet_text = (r.get("snippet", "") + " " + r.get("htmlSnippet", "")).lower()
        cagr_match = re.search(r'(\d+\.?\d*)\s*%\s*(?:cagr|compound annual growth rate)', snippet_text)
        if cagr_match: cagr_found = float(cagr_match.group(1))
        market_value_matches = re.findall(r'(\$?\d{1,3}(?:[,\.]\d{3})*(?:[,\.]\d+)?)\s*(?:usd|€)?\s*(billion|million|bn|m|billones|millones)?\s*(?:in|by)?\s*(\d{4})', snippet_text, re.IGNORECASE)
        for match in market_value_matches:
            value_in_millions = convertir_valor(match[0], match[1])
            if value_in_millions is not None: market_values.append((int(match[2]), value_in_millions))
        if cagr_found and market_values: break
    tam, sam, som, snippet = obtener_tam_sam_som_internal(product_query)
    return cagr_found, market_values, tam, sam, som, snippet

def obtener_tam_sam_som_internal(product_query):
    busqueda = f"{product_query} market size in USD total addressable serviceable obtainable"
    organic_results_raw = search_google_custom_search(busqueda)
    organic_results = organic_results_raw.get("items", [])
    tam_val, sam_val, som_val, source_snippet = None, None, None, ""
    keywords = {"TAM": ["total addressable market", "tam"], "SAM": ["serviceable available market", "sam"], "SOM": ["serviceable obtainable market", "som"]}
    for result in organic_results:
        text_to_search = (result.get("snippet", "") + " " + result.get("title", "")).lower()
        for market_type, kws in keywords.items():
            for kw in kws:
                if kw in text_to_search:
                    match = re.search(fr'(\$?\d{{1,3}}(?:[,\.]\d{{3}})*(?:[,\.]\d+)?)\s*(billion|million|bn|m|billones|millones)?', text_to_search, re.IGNORECASE)
                    if match:
                        value = convertir_valor(match.group(1), match.group(2))
                        if value is not None:
                            if market_type == "TAM" and tam_val is None: tam_val = value; source_snippet = result.get("snippet", "")
                            elif market_type == "SAM" and sam_val is None: sam_val = value
                            elif market_type == "SOM" and som_val is None: som_val = value
    return tam_val, sam_val, som_val, source_snippet

def extraer_empresas_relacionadas(product_query):
    busqueda = f"{product_query} competitors OR similar companies crunchbase.com"
    resultados_raw = search_google_custom_search(busqueda)
    empresas, seen_companies = [], set()
    for r in resultados_raw.get("items", [])[:15]:
        link, title, snippet = r.get("link", ""), r.get("title", ""), r.get("snippet", "")
        if "crunchbase.com" not in link.lower(): continue
        company_name = title.split(" - ")[0].strip()
        if not company_name or company_name in seen_companies: continue
        seen_companies.add(company_name)
        financing_match = re.search(r'(\$?\d+(?:[,\.]\d+)?)\s*([MB]n|million|billion|M|B)', snippet, re.IGNORECASE)
        financing = f"{financing_match.group(1)}{financing_match.group(2)[0].upper()}" if financing_match else "N/A"
        empresas.append({"Name": company_name, "Country": "N/A", "Financing": financing, "Approach": snippet[:150] + "..."})
    return empresas

def extraer_colaboraciones_ma(product_query):
    busqueda = f"{product_query} M&A collaborations OR acquisitions OR partnerships site:crunchbase.com"
    resultados_raw = search_google_custom_search(busqueda)
    collaborations, seen_links = [], set()
    for r in resultados_raw.get("items", [])[:10]:
        link = r.get("link", "")
        if link in seen_links: continue
        seen_links.add(link)
        snippet = r.get("snippet", "")
        year_match = re.search(r'\b(20\d{2})\b', snippet)
        year = year_match.group(1) if year_match else "N/A"
        collaborations.append({"Year": year, "Companies": r.get("title", ""), "Description of the agreement": snippet, "Type of agreement": "N/A"})
    return collaborations

def generar_drivers_barreras(product_query):
    search_queries = [f"{product_query} market drivers and barriers", f"{product_query} growth factors challenges", f"{product_query} industry trends opportunities threats"]
    context_text = ""
    for sq in search_queries:
        results_raw = search_google_custom_search(sq)
        for r in results_raw.get("items", [])[:5]:
            context_text += r.get("snippet", "") + " "
    if not context_text.strip(): return [], []
    prompt = f'Analiza el siguiente texto y extrae los "drivers" y "barreras" para el mercado de "{product_query}". Formatea como JSON con claves "drivers" y "barreras", donde cada uno es una lista de objetos. Cada driver tiene "tendency", "impact", "opportunity". Cada barrera tiene "difficulty", "tolerance", "limits".\nTexto: {context_text[:4000]}'
    
    payload = {
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {
            "responseMimeType": "application/json",
            "responseSchema": {
                "type": "OBJECT",
                "properties": {
                    "drivers": { "type": "ARRAY", "items": { "type": "OBJECT", "properties": { "tendency": {"type": "STRING"}, "impact": {"type": "STRING"}, "opportunity": {"type": "STRING"} } } },
                    "barriers": { "type": "ARRAY", "items": { "type": "OBJECT", "properties": { "difficulty": {"type": "STRING"}, "tolerance": {"type": "STRING"}, "limits": {"type": "STRING"} } } }
                }
            }
        }
    }
    
    drivers_list, barriers_list = [], []
    try:
        apiUrl = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash-latest:generateContent?key={GEMINI_API_KEY}"
        response = requests.post(apiUrl, headers={'Content-Type': 'application/json'}, data=json.dumps(payload), verify=False)
        response.raise_for_status()
        result = response.json()
        if result and result.get("candidates"):
            json_response_str = result["candidates"][0]["content"]["parts"][0]["text"]
            parsed_json = json.loads(json_response_str)
            drivers_list = parsed_json.get("drivers", [])
            barriers_list = parsed_json.get("barriers", [])
    except Exception as e:
        st.error(f"Error al generar drivers y barreras con Gemini: {e}")
        drivers_list = [{"tendency": "Increased demand for remote care", "impact": "Higher adoption of digital tools", "opportunity": "Offer 24/7 automated support"}]
        barriers_list = [{"difficulty": "Data privacy concerns (HIPAA)", "tolerance": "User trust in AI for diagnosis", "limits": "Regulatory hurdles"}]
    
    return drivers_list, barriers_list

# --- 4. FUNCIONES DE EXPORTACIÓN (PDF Y PPTX) ---
class PDF(FPDF):
    def header(self):
        try: self.image("vicomtech_logo.png", 10, 8, 33)
        except FileNotFoundError: self.cell(0, 10, 'Vicomtech Logo', 0, 0, 'L')
        self.set_font('Arial', 'B', 15)
        self.cell(80); self.cell(30, 10, 'Market Potential Report', 0, 0, 'C'); self.ln(20)
    def footer(self):
        self.set_y(-15); self.set_font('Arial', 'I', 8); self.cell(0, 10, f'Página {self.page_no()}', 0, 0, 'C')
    def chapter_title(self, title):
        self.set_font('Arial', 'B', 12); self.set_fill_color(220, 220, 220); self.cell(0, 6, title, 0, 1, 'L', 1); self.ln(4)
    def chapter_body(self, body):
        self.set_font('Arial', '', 11)
        self.multi_cell(0, 5, body.encode('latin-1', 'replace').decode('latin-1'))
        self.ln()
    
    def add_dataframe_to_pdf(self, df):
        # **FIXED**: Re-implemented table drawing logic to be more robust
        if df.empty:
            self.chapter_body("No data available.")
            return

        self.set_font('Arial', 'B', 9)
        line_height = self.font_size * 2
        
        available_width = self.w - self.l_margin - self.r_margin
        if "Approach" in df.columns: col_widths = [available_width * 0.2, available_width * 0.15, available_width * 0.15, available_width * 0.5]
        elif "Description of the agreement" in df.columns: col_widths = [available_width * 0.1, available_width * 0.3, available_width * 0.4, available_width*0.2]
        else: col_widths = [available_width / len(df.columns)] * len(df.columns)
        
        # Header
        for i, col_name in enumerate(df.columns):
            self.cell(col_widths[i], line_height, str(col_name), 1, 0, 'C')
        self.ln()

        # Data rows
        self.set_font('Arial', '', 8)
        for _, row in df.iterrows():
            y_before_row = self.get_y()
            x_pos = self.l_margin
            max_y = y_before_row

            # First pass: Write text to find max row height
            for i, item in enumerate(row):
                self.set_xy(x_pos, y_before_row)
                text = str(item).encode('latin-1', 'replace').decode('latin-1')
                self.multi_cell(col_widths[i], 5, text, border=0, align='L')
                if self.get_y() > max_y:
                    max_y = self.get_y()
                x_pos += col_widths[i]

            row_height = max_y - y_before_row
            
            # Second pass: Draw borders and text
            self.set_y(y_before_row) # Go back to the start of the row
            x_pos = self.l_margin
            for i, item in enumerate(row):
                self.set_xy(x_pos, y_before_row)
                self.rect(x_pos, y_before_row, col_widths[i], row_height)
                text = str(item).encode('latin-1', 'replace').decode('latin-1')
                self.multi_cell(col_widths[i], 5, text, border=0, align='L')
                x_pos += col_widths[i]
            
            self.set_y(y_before_row + row_height)


def create_report_pdf(query, data):
    pdf = PDF()
    pdf.add_page()
    pdf.set_font('Arial', 'B', 16)
    pdf.cell(0, 10, f'Análisis de Mercado para: {query}', 0, 1, 'C'); pdf.ln(10)
    pdf.chapter_title('1. Colección de Números Cruciales')
    pdf.chapter_body(f"CAGR: {data.get('cagr', 'N/A') or 'N/A'}%")
    if data.get('market_values'):
        pdf.chapter_body("Valores de Mercado por Año:")
        for year, value in sorted(data['market_values'], key=lambda x: x[0]): pdf.chapter_body(f"  - {year}: {value} M USD")
    pdf.chapter_title('2. TAM, SAM & SOM')
    pdf.chapter_body(f"TAM: {data.get('tam', 'N/A')} M USD\nSAM: {data.get('sam', 'N/A')} M USD\nSOM: {data.get('som', 'N/A')} M USD")
    pdf.add_page()
    pdf.chapter_title('3. Competitive Landscape')
    if data.get('empresas_competidoras'): pdf.add_dataframe_to_pdf(pd.DataFrame(data['empresas_competidoras']))
    pdf.chapter_title('4. Drivers & Barriers')
    if data.get('drivers'):
        for d in data['drivers']: pdf.chapter_body(f"Driver: {d.get('tendency', 'N/A')}")
    if data.get('barriers'):
        for b in data['barriers']: pdf.chapter_body(f"Barrera: {b.get('difficulty', 'N/A')}")
    pdf.add_page()
    pdf.chapter_title('5. M&A Collaborations')
    if data.get('ma_collaborations'): pdf.add_dataframe_to_pdf(pd.DataFrame(data['ma_collaborations']))
    return pdf.output(dest='S').encode('latin-1')

def add_table_to_slide(slide, df, left, top, width, height):
    shape = slide.shapes.add_table(df.shape[0] + 1, df.shape[1], left, top, width, height)
    table = shape.table
    for col_i, col_name in enumerate(df.columns): table.cell(0, col_i).text = str(col_name)
    for row_i, row_data in df.iterrows():
        for col_i, cell_data in enumerate(row_data): table.cell(row_i + 1, col_i).text = str(cell_data)

def create_report_pptx(query, data):
    template_path = 'vicomtech_template.pptx'
    try: prs = Presentation(template_path)
    except Exception: prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    table_slide_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = f"Market Potential Report: {query}"
    if len(slide.placeholders) > 1: slide.placeholders[1].text = "Generado por Market Intelligence Vicomtech"
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "1. & 2. Crucial Numbers & Market Sizing"
    tf = slide.placeholders[1].text_frame; tf.clear()
    p = tf.add_paragraph(); p.text = f"CAGR: {data.get('cagr', 'N/A') or 'N/A'}%"; p.level = 0
    if data.get('market_values'):
        p = tf.add_paragraph(); p.text = "Market Values:"; p.level = 0
        for year, value in sorted(data['market_values'], key=lambda x: x[0]):
             p = tf.add_paragraph(); p.text = f"{year}: {value} M USD"; p.level = 1
    p = tf.add_paragraph(); p.text = f"TAM/SAM/SOM: {data.get('tam', 'N/A')}M / {data.get('sam', 'N/A')}M / {data.get('som', 'N/A')}M USD"; p.level = 0

    if data.get('empresas_competidoras'):
        slide = prs.slides.add_slide(table_slide_layout)
        slide.shapes.title.text = "3. Competitive Landscape"
        add_table_to_slide(slide, pd.DataFrame(data['empresas_competidoras']), Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "4. Drivers & Barriers"
    tf = slide.placeholders[1].text_frame; tf.clear()
    p = tf.add_paragraph(); p.text = "Drivers"; p.level = 0
    if data.get('drivers'):
        for d in data['drivers']: p = tf.add_paragraph(); p.text = d.get('tendency', 'N/A'); p.level = 1
    p = tf.add_paragraph(); p.text = "Barriers"; p.level = 0
    if data.get('barriers'):
        for b in data['barriers']: p = tf.add_paragraph(); p.text = b.get('difficulty', 'N/A'); p.level = 1

    if data.get('ma_collaborations'):
        slide = prs.slides.add_slide(table_slide_layout)
        slide.shapes.title.text = "5. M&A Collaborations"
        add_table_to_slide(slide, pd.DataFrame(data['ma_collaborations']), Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()

# --- 5. LÓGICA PRINCIPAL DE LA APLICACIÓN ---
try:
    st.image("vicomtech_logo.png", width=300)
except Exception:
    st.warning("No se encontró el archivo 'vicomtech_logo.png'.")

st.title("Market Potential Generator")
st.markdown("**¡Bienvenido!** Introduce el nombre de un producto o solución para generar un informe preliminar.")
query = st.text_input("Nombre del producto o solución", key="main_product_query_input")

if query:
    st.markdown(f"<h2 style='color:#000000;'>Preliminary Market Overview for: <i>{query}</i></h2>", unsafe_allow_html=True)
    analysis_data = {}
    with st.spinner("Realizando análisis de mercado completo..."):
        analysis_data = run_full_analysis(query)
    st.success("Análisis completado.")

    # --- Mostrar todos los resultados detallados en la UI ---
    st.markdown("### Informes por Fuente")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("**Statista**")
        if analysis_data.get('statista_reports'):
            for titulo, link in analysis_data['statista_reports']: st.markdown(f"- [{titulo}]({link})")
    with col2:
        st.markdown("**Market Data Forecast**")
        if analysis_data.get('market_data_reports'):
            for titulo, link in analysis_data['market_data_reports']: st.markdown(f"- [{titulo}]({link})")
    with col3:
        st.markdown("**Precedence Research**")
        if analysis_data.get('precedence_reports'):
            for titulo, link in analysis_data['precedence_reports']: st.markdown(f"- [{titulo}]({link})")

    st.markdown("---")
    st.markdown("<h2 style='color:#000000;'>1. Colección de Números Cruciales</h2>", unsafe_allow_html=True)
    st.markdown(f"**CAGR:** {analysis_data.get('cagr')}%" if analysis_data.get('cagr') else "**CAGR:** No encontrado")
    if analysis_data.get('market_values'):
        st.markdown("**Valores de Mercado por Año:**")
        for year, value in sorted(analysis_data['market_values'], key=lambda x: x[0]): st.markdown(f"  - **{year}:** {value} millones USD")
    
    st.markdown("---")
    st.markdown("<h2 style='color:#000000;'>2. TAM, SAM y SOM</h2>", unsafe_allow_html=True)
    if analysis_data.get('tam_sam_som_snippet'): st.markdown(f"> _{analysis_data['tam_sam_som_snippet']}_")
    st.write(f"**TAM:** {analysis_data.get('tam') or 'N/A'} M USD | **SAM:** {analysis_data.get('sam') or 'N/A'} M USD | **SOM:** {analysis_data.get('som') or 'N/A'} M USD")
    
    st.markdown("---")
    st.markdown("<h2 style='color:#000000;'>3. Competitive Landscape</h2>", unsafe_allow_html=True)
    if not pd.DataFrame(analysis_data.get('empresas_competidoras')).empty:
        st.dataframe(pd.DataFrame(analysis_data['empresas_competidoras']), use_container_width=True)
    else:
        st.warning("No se encontraron empresas relevantes.")

    st.markdown("---")
    st.markdown("<h2 style='color:#000000;'>4. Drivers y Barreras</h2>", unsafe_allow_html=True)
    st.markdown("### Drivers (Factores Impulsores)")
    if analysis_data.get('drivers'):
        for d in analysis_data['drivers']: st.markdown(f"- **{d.get('tendency','')}**: {d.get('impact','')}. **Oportunidad:** {d.get('opportunity','N/A')}")
    else: st.warning("No se pudieron generar drivers.")
    
    st.markdown("### Barreras (Desafíos)")
    if analysis_data.get('barriers'):
        for b in analysis_data['barriers']: st.markdown(f"- **{b.get('difficulty','')}**: {b.get('tolerance','')}. **Límites:** {b.get('limits','N/A')}")
    else: st.warning("No se pudieron generar barreras.")

    st.markdown("---")
    st.markdown("<h2 style='color:#000000;'>5. M&A Collaborations</h2>", unsafe_allow_html=True)
    if not pd.DataFrame(analysis_data.get('ma_collaborations')).empty:
        st.dataframe(pd.DataFrame(analysis_data['ma_collaborations']), use_container_width=True)
    else:
        st.warning("No se encontraron colaboraciones de M&A.")

    # --- 6. BOTONES DE DESCARGA (PDF Y PPTX) ---
    st.markdown("---")
    st.markdown("<h2 style='color:#000000;'>Descargar Informe Completo</h2>", unsafe_allow_html=True)
    
    dl_col1, dl_col2 = st.columns(2)
    with dl_col1:
        with st.spinner("Generando PDF..."):
            pdf_bytes = create_report_pdf(query, analysis_data)
            st.download_button(label="📥 Descargar Reporte en PDF", data=pdf_bytes, file_name=f"Reporte_Mercado_{query.replace(' ', '_')}.pdf", mime="application/pdf")
    with dl_col2:
        with st.spinner("Generando PPTX..."):
            pptx_bytes = create_report_pptx(query, analysis_data)
            st.download_button(label="📥 Descargar Presentación en PPTX", data=pptx_bytes, file_name=f"Presentacion_Mercado_{query.replace(' ', '_')}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    
    st.markdown("---")
    st.info("Para un análisis más profundo, considere la investigación manual de las fuentes.")
